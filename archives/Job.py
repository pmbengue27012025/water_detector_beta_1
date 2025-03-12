from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Créer une présentation
prs = Presentation()

# ---- Diapositive 1 (Titre) ----
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Surveillance Environnementale avec Sentinel-2 et Open-Meteo"
subtitle.text = "Par [Votre Nom]\n[Date]"

# ---- Diapositive 2 (Introduction) ----
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Introduction"
content.text = (
    "Objectif :\n"
    "- Analyser NDVI, NDWI/MNDWI, données météo\n"
    "- Public cible : Agronomes, écologues, gestionnaires"
)

# ---- Diapositive 3 (Technologies) ----
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Technologies Utilisées"
content.text = (
    "Backend :\n"
    "- FastAPI, Sentinel Hub, Open-Meteo\n"
    "Traitement d'images :\n"
    "- Rasterio, NumPy"
)

# ---- Diapositive 4 (Fonctionnalités) ----
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Fonctionnalités Clés"
content.text = (
    "1. Analyse TIFF : NDVI, NDWI, détection d'eau\n"
    "2. Données météo : Températures, précipitations\n"
    "3. API : /analyze-tiff et /get-analysis"
)

# ---- Diapositive 5 (Architecture) ----
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Architecture"
content.text = (
    "Workflow :\n"
    "Utilisateur → FastAPI → Traitement → Résultats JSON\n"
    "Intégration : Sentinel-2 + Open-Meteo"
)

# ---- Continuez pour les autres diapositives... ----

# Sauvegarder le fichier
prs.save("Analyse_Environnementale.pptx")
print("Présentation générée avec succès !")